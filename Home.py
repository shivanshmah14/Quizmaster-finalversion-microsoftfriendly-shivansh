import streamlit as st
from pathlib import Path
import database as db
import time
import requests
import json
import os

st.set_page_config(
    page_title="QuizMaster - Home",
    page_icon="🎯",
    layout="wide"
)

# ── API config ─────────────────────────────────────────────────────────────────
SARVAM_API_KEY = os.getenv("SARVAM_API_KEY", "sk_02syy9dv_Kl4BUcngJpzdwyhoYs0Tgk68")
SARVAM_CHAT_URL = "https://api.sarvam.ai/v1/chat/completions"


def initialize_session_state():
    defaults = {
        "logged_in": False,
        "user_id": None,
        "username": None,
        "is_admin": False,
        "player_name": "",
        "score": 0,
        "current_question": 0,
        "game_active": False,
        "selected_category": None,
        "answers_given": [],
        "correct_answers": 0,
        "time_remaining": 30,
        "questions": [],
        "answered_current": False,
        "show_result": False,
        "timer_start": None,
        "timer_expired": False,
        "shiva_open": False,
        "shiva_messages": [],
        "file_quiz_questions": [],
        "file_quiz_index": 0,
        "file_quiz_score": 0,
        "file_quiz_active": False,
        "file_quiz_answered": False,
        "file_quiz_last_answer": None,
        "show_file_quiz_result": False,
    }
    for key, value in defaults.items():
        if key not in st.session_state:
            st.session_state[key] = value


def call_shiva_ai(user_message, history):
    system_prompt = (
        "You are Shiva AI, a quiz study assistant. Follow these rules strictly:\n"
        "1. ALWAYS respond in English only, no matter what language the user writes in.\n"
        "2. NEVER show your thinking, reasoning, or internal monologue. Output the final answer only.\n"
        "3. Keep responses to 1-3 sentences. Be direct and concise.\n"
        "4. Help with quiz topics, studying, and learning only."
    )
    messages = [{"role": "system", "content": system_prompt}]
    for msg in history[-6:]:
        messages.append(msg)
    messages.append({"role": "user", "content": user_message})
    headers = {
        "Authorization": f"Bearer {SARVAM_API_KEY}",
        "Content-Type": "application/json"
    }
    payload = {"model": "sarvam-m", "messages": messages, "max_tokens": 150}
    try:
        resp = requests.post(SARVAM_CHAT_URL, headers=headers, json=payload, timeout=15)
        resp.raise_for_status()
        return resp.json()["choices"][0]["message"]["content"].strip()
    except requests.exceptions.Timeout:
        return "Request timed out. Please try again."
    except requests.exceptions.HTTPError as e:
        code = e.response.status_code if e.response else "?"
        if code == 401:
            return "API key error - please update your Sarvam AI key."
        elif code == 429:
            return "Too many requests - wait a moment and try again."
        return f"API error ({code}). Please try again."
    except Exception as e:
        return f"Error: {str(e)}"


def generate_quiz_from_text(text):
    headers = {
        "Authorization": f"Bearer {SARVAM_API_KEY}",
        "Content-Type": "application/json"
    }
    prompt = (
        "You are a quiz creator. Read the text below and generate exactly 5 multiple-choice questions.\n\n"
        "CRITICAL: Return ONLY a raw JSON array. No explanation, no markdown, no code fences.\n"
        "Start your entire response with [ and end with ]. Nothing before or after.\n\n"
        "Each object must have:\n"
        "- \"question\": string\n"
        "- \"options\": array of exactly 4 strings\n"
        "- \"correct\": integer 0-3 (index of correct option)\n"
        "- \"explanation\": string (one sentence)\n\n"
        f"Text:\n{text[:3000]}"
    )
    payload = {
        "model": "sarvam-m",
        "messages": [{"role": "user", "content": prompt}],
        "max_tokens": 1500
    }
    try:
        resp = requests.post(SARVAM_CHAT_URL, headers=headers, json=payload, timeout=30)
        resp.raise_for_status()
        raw = resp.json()["choices"][0]["message"]["content"].strip()
        # Strip markdown fences if the model adds them
        if "```" in raw:
            parts = raw.split("```")
            for part in parts:
                part = part.strip()
                if part.startswith("json"):
                    part = part[4:].strip()
                if part.startswith("["):
                    raw = part
                    break
        # Extract just the JSON array
        start = raw.find("[")
        end = raw.rfind("]")
        if start != -1 and end != -1:
            raw = raw[start:end+1]
        return json.loads(raw)
    except Exception:
        return []


def extract_text_from_file(uploaded_file):
    name = uploaded_file.name.lower()
    plain_text_extensions = (
        ".txt", ".md", ".js", ".jsx", ".ts", ".tsx",
        ".py", ".css", ".html", ".htm", ".json", ".csv", ".xml"
    )
    try:
        if any(name.endswith(ext) for ext in plain_text_extensions):
            return uploaded_file.read().decode("utf-8", errors="ignore")

        elif name.endswith(".pdf"):
            from pypdf import PdfReader
            import io
            reader = PdfReader(io.BytesIO(uploaded_file.read()))
            text = "\n".join(page.extract_text() or "" for page in reader.pages)
            if not text.strip():
                st.warning("PDF appears to be scanned/image-based. Text extraction may be limited.")
            return text

        elif name.endswith(".docx"):
            from docx import Document
            import io
            doc = Document(io.BytesIO(uploaded_file.read()))
            return "\n".join(p.text for p in doc.paragraphs if p.text.strip())

        elif name.endswith(".pptx"):
            from pptx import Presentation
            import io
            prs = Presentation(io.BytesIO(uploaded_file.read()))
            texts = []
            for slide in prs.slides:
                for shape in slide.shapes:
                    if shape.has_text_frame:
                        texts.append(shape.text_frame.text)
            return "\n".join(texts)

        else:
            st.error(f"File type not supported: {name.split('.')[-1].upper()}")
            return ""

    except ImportError as e:
        st.error(f"Missing library: {e}. Run: pip install pypdf python-docx python-pptx")
        return ""
    except Exception as e:
        st.error(f"Could not read file: {e}")
        return ""


def show_login():
    st.subheader("Login")
    with st.form("login_form"):
        username = st.text_input("Username")
        password = st.text_input("Password", type="password")
        submit = st.form_submit_button("Login", use_container_width=True)
        if submit:
            if username and password:
                user = db.verify_user(username, password)
                if user:
                    st.session_state.logged_in = True
                    st.session_state.user_id = user['id']
                    st.session_state.username = user['username']
                    st.session_state.is_admin = user['is_admin']
                    st.session_state.player_name = user['username']
                    st.success(f"Welcome back, {username}!")
                    st.rerun()
                else:
                    st.error("Invalid username or password")
            else:
                st.warning("Please fill in all fields")


def show_register():
    st.subheader("Create Account")
    with st.form("register_form"):
        username = st.text_input("Username")
        password = st.text_input("Password", type="password", help="Minimum 6 characters")
        password_confirm = st.text_input("Confirm Password", type="password")
        submit = st.form_submit_button("Register", use_container_width=True)
        if submit:
            if not username or not password:
                st.error("Please fill in all fields")
            elif password != password_confirm:
                st.error("Passwords do not match")
            else:
                success, message = db.create_user(username, password)
                if success:
                    st.success(message)
                    st.info("You can now login with your credentials")
                else:
                    st.error(message)


def show_auth_page():
    st.markdown("""
        <div style="text-align:center;padding:2rem 0;
        background:linear-gradient(135deg,#667eea 0%,#764ba2 100%);
        border-radius:10px;margin-bottom:2rem;color:white;">
            <h1>🎯 QuizMaster</h1>
            <p>Test Your Knowledge and Challenge Yourself!</p>
        </div>
    """, unsafe_allow_html=True)
    st.markdown("### Welcome! Please login or create an account to continue.")
    tab1, tab2 = st.tabs(["Login", "Register"])
    with tab1:
        show_login()
    with tab2:
        show_register()


def show_file_quiz_section():
    st.markdown("---")
    st.markdown("### 🤖 Generate a Quiz from Your File")
    st.caption("Upload any file — Shiva AI will create a quiz from it!")

    uploaded_file = st.file_uploader(
        "Upload a file",
        type=[
            "txt", "md", "pdf", "docx", "pptx",
            "js", "jsx", "ts", "tsx", "html", "htm", "css",
            "py", "json", "csv", "xml"
        ],
        key="quiz_file_uploader"
    )

    if uploaded_file is not None:
        file_text = extract_text_from_file(uploaded_file)
        if not file_text.strip():
            st.error("Could not extract text from this file.")
            return
        word_count = len(file_text.split())
        st.success(f"✅ File loaded: {uploaded_file.name} ({word_count} words)")

        if st.button("🤖 Generate Quiz from File", use_container_width=True, type="primary"):
            with st.spinner("Shiva AI is reading your file and creating questions..."):
                questions = generate_quiz_from_text(file_text)
            if questions:
                st.session_state.file_quiz_questions = questions
                st.session_state.file_quiz_index = 0
                st.session_state.file_quiz_score = 0
                st.session_state.file_quiz_active = True
                st.session_state.file_quiz_answered = False
                st.session_state.file_quiz_last_answer = None
                st.session_state.show_file_quiz_result = False
                st.rerun()
            else:
                st.error("❌ Could not generate questions. Try a file with more text content.")

    if st.session_state.get("file_quiz_active") and not st.session_state.get("show_file_quiz_result"):
        questions = st.session_state.file_quiz_questions
        idx = st.session_state.file_quiz_index
        total = len(questions)

        if idx >= total:
            st.session_state.file_quiz_active = False
            st.session_state.show_file_quiz_result = True
            st.rerun()
            return

        q = questions[idx]
        st.markdown("---")
        st.markdown(f"#### 🤖 Shiva AI Quiz — Question {idx + 1} of {total}")
        st.progress(idx / total)
        st.markdown(f"""
            <div style="padding:1.5rem;border-radius:10px;background:#f8f9fa;
            border:2px solid #667eea;margin:1rem 0;">
                <h3>{q['question']}</h3>
            </div>
        """, unsafe_allow_html=True)

        if not st.session_state.file_quiz_answered:
            for i, option in enumerate(q['options']):
                if st.button(f"{chr(65+i)}) {option}", key=f"fq_{idx}_{i}", use_container_width=True):
                    st.session_state.file_quiz_last_answer = i
                    st.session_state.file_quiz_answered = True
                    if i == q['correct']:
                        st.session_state.file_quiz_score += 1
                    st.rerun()
        else:
            last = st.session_state.file_quiz_last_answer
            for i, option in enumerate(q['options']):
                if i == q['correct']:
                    st.success(f"{chr(65+i)}) {option}")
                elif i == last:
                    st.error(f"{chr(65+i)}) {option}")
                else:
                    st.write(f"{chr(65+i)}) {option}")

            if last == q['correct']:
                st.success("✅ Correct!")
            else:
                st.error(f"❌ Wrong! Correct answer: {q['options'][q['correct']]}")
            if q.get('explanation'):
                st.info(f"💡 {q['explanation']}")

            col1, col2, col3 = st.columns([1, 2, 1])
            with col2:
                label = "Next Question ➡️" if idx < total - 1 else "Finish Quiz 🏁"
                if st.button(label, use_container_width=True, type="primary"):
                    st.session_state.file_quiz_index += 1
                    st.session_state.file_quiz_answered = False
                    st.session_state.file_quiz_last_answer = None
                    if idx + 1 >= total:
                        st.session_state.file_quiz_active = False
                        st.session_state.show_file_quiz_result = True
                    st.rerun()

    if st.session_state.get("show_file_quiz_result"):
        total = len(st.session_state.file_quiz_questions)
        score = st.session_state.file_quiz_score
        pct = (score / total * 100) if total > 0 else 0
        st.markdown("---")
        st.markdown("#### 🤖 Shiva AI Quiz — Results")
        col1, col2 = st.columns(2)
        with col1:
            st.metric("Score", f"{score}/{total}")
        with col2:
            st.metric("Accuracy", f"{pct:.0f}%")
        if pct >= 80:
            st.success("🌟 Excellent work!")
        elif pct >= 50:
            st.info("👍 Good effort! Keep studying.")
        else:
            st.warning("💪 Keep studying — you've got this!")
        if st.button("Try Another File", use_container_width=True):
            st.session_state.show_file_quiz_result = False
            st.session_state.file_quiz_questions = []
            st.rerun()


def reset_timer():
    st.session_state.timer_start = None
    st.session_state.timer_expired = False


def show_quiz_home():
    st.markdown("""
        <style>
        .main-header {
            text-align:center; padding:2rem 0;
            background:linear-gradient(135deg,#667eea 0%,#764ba2 100%);
            border-radius:10px; margin-bottom:2rem; color:white;
        }
        .category-card {
            padding:1.5rem; border-radius:10px;
            background:#f8f9fa; border-left:5px solid #667eea; margin:1rem 0;
        }
        </style>
    """, unsafe_allow_html=True)

    st.markdown("""
        <div class="main-header">
            <h1>🎯 QuizMaster</h1>
            <p>Test Your Knowledge and Challenge Yourself!</p>
        </div>
    """, unsafe_allow_html=True)

    with st.sidebar:
        st.success(f"✅ Logged in as: **{st.session_state.username}**")
        if st.session_state.is_admin:
            st.info("🔑 Admin Account")
            if st.button("⚙️ Admin Panel", use_container_width=True):
                st.switch_page("pages/4_Admin.py")
        st.markdown("---")
        if st.button("🚪 Logout", use_container_width=True):
            for key in list(st.session_state.keys()):
                del st.session_state[key]
            st.rerun()

        st.markdown("---")
        st.markdown("### 🤖 Shiva AI")
        st.caption("Ask anything about your quiz topics!")

        chat_container = st.container(height=260)
        with chat_container:
            if not st.session_state.shiva_messages:
                st.markdown(
                    "<div style='color:#888;font-size:13px;font-style:italic;'>"
                    "Hi! I'm Shiva AI. Ask me anything about your quiz topics.</div>",
                    unsafe_allow_html=True
                )
            for msg in st.session_state.shiva_messages:
                role_label = "You" if msg["role"] == "user" else "Shiva AI"
                bg = "#e0e7ff" if msg["role"] == "user" else "#f1f5f9"
                align = "right" if msg["role"] == "user" else "left"
                st.markdown(
                    f"<div style='text-align:{align};margin-bottom:6px;'>"
                    f"<span style='font-size:11px;color:#888;'>{role_label}</span><br>"
                    f"<span style='background:{bg};padding:6px 10px;border-radius:10px;"
                    f"font-size:13px;display:inline-block;max-width:90%;text-align:left;'>"
                    f"{msg['content']}</span></div>",
                    unsafe_allow_html=True
                )

        user_input = st.chat_input("Ask Shiva AI...", key="shiva_chat_input")
        if user_input:
            st.session_state.shiva_messages.append({"role": "user", "content": user_input})
            history = [m for m in st.session_state.shiva_messages if m["role"] != "system"]
            with st.spinner("Shiva AI is thinking..."):
                reply = call_shiva_ai(user_input, history[:-1])
            st.session_state.shiva_messages.append({"role": "assistant", "content": reply})
            st.rerun()

        if st.button("🗑️ Clear Chat", use_container_width=True, key="clear_shiva"):
            st.session_state.shiva_messages = []
            st.rerun()

    st.markdown("### 👋 Welcome to QuizMaster!")
    st.write("An interactive learning app to test your knowledge across multiple categories.")

    categories_list = db.get_all_categories()

    if not categories_list:
        st.warning("⚠️ No quiz categories available yet.")
        if st.session_state.is_admin:
            st.info("As an admin, you can add questions in the Admin Panel.")
        else:
            st.info("Please ask an admin to add quiz questions.")
    else:
        st.markdown("### 📚 Choose a Category")
        with st.expander("⚙️ Quiz Settings"):
            timer_seconds = st.slider("Seconds per question", min_value=10, max_value=120, value=30, step=5)
            st.session_state.quiz_timer_seconds = timer_seconds

        cols = st.columns(2)
        for idx, category_name in enumerate(categories_list):
            questions = db.get_questions_by_category(category_name)
            with cols[idx % 2]:
                st.markdown(f"""
                    <div class="category-card">
                        <h4>{category_name}</h4>
                        <p>{len(questions)} questions available</p>
                    </div>
                """, unsafe_allow_html=True)
                if st.button(f"▶️ Start {category_name} Quiz", key=f"start_{category_name}", use_container_width=True):
                    st.session_state.selected_category = category_name
                    st.session_state.current_question = 0
                    st.session_state.score = 0
                    st.session_state.game_active = True
                    st.session_state.answers_given = []
                    st.session_state.correct_answers = 0
                    st.session_state.answered_current = False
                    st.session_state.show_result = False
                    reset_timer()
                    st.switch_page("pages/1_Quiz.py")

    st.markdown("---")
    st.markdown("### 📊 Quick Stats")
    col1, col2, col3 = st.columns(3)
    with col1:
        st.metric("Categories", len(categories_list) if categories_list else 0)
    with col2:
        total_q = sum(len(db.get_questions_by_category(cat)) for cat in categories_list) if categories_list else 0
        st.metric("Total Questions", total_q)
    with col3:
        user_stats = db.get_user_statistics(st.session_state.user_id)
        games_played = user_stats['total_games'] if user_stats else 0
        st.metric("Games Played", games_played)

    show_file_quiz_section()

    st.markdown("---")
    st.markdown("<div style='text-align:center;color:#666;'><p>Good luck with your quiz! 🍀</p></div>",
                unsafe_allow_html=True)


def main():
    initialize_session_state()
    if not st.session_state.logged_in:
        show_auth_page()
    else:
        show_quiz_home()

main()